from flask import Blueprint, jsonify, request
import google.generativeai as genai
from dotenv import load_dotenv
from pptx import Presentation
from models.summary_model import summaries_collection
from models.mcq_model import mcqs_collection
from models.quiz_model import quiz_collection
from models.note_model import notes_collection
from utils.db import db
from utils.paths import resolve_note_filepath, resolve_upload_filepath
from utils.errors import api_error, handle_exception
from datetime import datetime
from datetime import timedelta
import fitz
import json
import os
import tempfile
import urllib.request

from routes.library_routes import track_library_file_event_by_filename

load_dotenv()

# ==========================================
# GEMINI CONFIG
# ==========================================

genai.configure(
    api_key=os.getenv("GEMINI_API_KEY")
)

model = genai.GenerativeModel("gemini-2.5-flash")

ai = Blueprint("ai", __name__)

chat_collection = db["chat_messages"]

# ==========================================
# HELPER FUNCTIONS
# ==========================================

def extract_text(filepath):

    if not filepath or not os.path.exists(filepath):
        raise FileNotFoundError("File not found")

    lower = filepath.lower()
    text = ""

    if lower.endswith(".pptx") or lower.endswith(".ppt"):
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    elif lower.endswith(".pdf"):
        pdf = fitz.open(filepath)
        for page in pdf:
            text += page.get_text()
        pdf.close()

    elif lower.endswith(".txt"):
        with open(filepath, "r", encoding="utf-8", errors="replace") as handle:
            text = handle.read()

    elif lower.endswith(".docx"):
        try:
            from docx import Document
            doc = Document(filepath)
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except ImportError:
            raise ValueError("DOCX support is not available on this server")

    else:
        raise ValueError("Unsupported file format. Use PDF, PPTX, DOCX, or TXT.")

    return text


def clean_json(raw_text):

    raw_text = raw_text.strip()
    raw_text = raw_text.replace("```json", "")
    raw_text = raw_text.replace("```", "")

    start = raw_text.find("[")
    end = raw_text.rfind("]")
    if start != -1 and end != -1 and end > start:
        return raw_text[start : end + 1].strip()

    start = raw_text.find("{")
    end = raw_text.rfind("}")
    if start != -1 and end != -1 and end > start:
        return raw_text[start : end + 1].strip()

    return raw_text.strip()


def parse_mcqs(raw_text):
    parsed = json.loads(clean_json(raw_text))

    if isinstance(parsed, list):
        return parsed

    if isinstance(parsed, dict):
        for key in ("mcqs", "questions", "items"):
            value = parsed.get(key)
            if isinstance(value, list):
                return value

    raise ValueError("AI response did not contain MCQs")


def study_text_for_filename(filename):
    summary_doc = summaries_collection.find_one({"filename": filename})

    if summary_doc:
        parts = summary_doc.get("summary", [])
        if isinstance(parts, list):
            text = "\n".join(str(p) for p in parts if p)
        else:
            text = str(parts or "")
        if text.strip():
            return text, summary_doc

    note = notes_collection.find_one({"filename": filename})
    filepath, temp_path = material_filepath(note, filename)

    if not os.path.exists(filepath):
        return None, summary_doc

    try:
        text = extract_text(filepath)
        return text[:15000], summary_doc
    finally:
        cleanup_temp_file(temp_path)


def material_filepath(note, filename):
    filepath = resolve_note_filepath(note) if note else resolve_upload_filepath(filename)
    if os.path.exists(filepath):
        return filepath, None

    cloudinary_url = (note or {}).get("cloudinary_url", "")
    if not cloudinary_url:
        return filepath, None

    original_name = (note or {}).get("original_name") or filename
    _, ext = os.path.splitext(original_name)
    fd, temp_path = tempfile.mkstemp(suffix=ext or ".pdf")
    os.close(fd)
    urllib.request.urlretrieve(cloudinary_url, temp_path)
    return temp_path, temp_path


def cleanup_temp_file(temp_path):
    if temp_path and os.path.exists(temp_path):
        os.remove(temp_path)


# ==========================================
# TEST ROUTE
# ==========================================

@ai.route("/test", methods=["GET"])
def test():

    return jsonify({
        "message": "AI Routes Working Successfully 🚀"
    })


# ==========================================
# SUMMARY ROUTE
# ==========================================

@ai.route("/summary/<path:filename>", methods=["GET"])
def generate_summary(filename):

    try:
        regenerate = request.args.get("regenerate") == "true"

        # -----------------------
        # Check MongoDB Cache
        # -----------------------

        existing = summaries_collection.find_one({
            "filename": filename
        })

        if existing and not regenerate:

            return jsonify({

                "filename": filename,

                "subject": existing.get("subject"),

                "unit": existing.get("unit"),

                "summary": existing["summary"],

                "key_points": existing.get("key_points", []),

                "important_topics": existing.get(
                    "important_topics",
                    []
                ),

                "source": "cache"

            })

        if regenerate:
            summaries_collection.delete_many({
                "filename": filename
            })

        # -----------------------
        # Get Subject & Unit
        # -----------------------

        note = notes_collection.find_one({
            "filename": filename
        })

        subject = ""
        unit = ""

        if note:
            subject = note.get("subject", "")
            unit = note.get("unit", "")

        filepath, temp_path = material_filepath(note, filename)

        if not os.path.exists(filepath):
            return api_error("File not found", 404)

        # -----------------------
        # Extract Text
        # -----------------------

        try:
            text = extract_text(filepath)
        finally:
            cleanup_temp_file(temp_path)

        if len(text.strip()) == 0:
            return api_error("No text found in this file", 400)

        # -----------------------
        # Gemini Prompt
        # -----------------------

        prompt = f"""
You are an expert teacher.

Analyze the study material.

Return ONLY valid JSON.

Format:

{{
    "summary":[
        "Paragraph 1",
        "Paragraph 2"
    ],

    "key_points":[
        "Point 1",
        "Point 2",
        "Point 3",
        "Point 4",
        "Point 5"
    ],

    "important_topics":[
        "Topic 1",
        "Topic 2",
        "Topic 3",
        "Topic 4",
        "Topic 5"
    ]
}}

Rules:

1. Do NOT return markdown.
2. Do NOT return explanation.
3. Return JSON only.
4. Summary should be student-friendly.
5. Keep important topics short.

Study Material:

{text[:15000]}
"""

        response = model.generate_content(prompt)

        raw = clean_json(response.text)
        summary_json = json.loads(raw)

        summary = summary_json.get("summary", [])

        key_points = summary_json.get("key_points", [])

        important_topics = summary_json.get(
            "important_topics",
            []
        )

        summaries_collection.insert_one({

            "filename": filename,

            "subject": subject,

            "unit": unit,

            "summary": summary,

            "key_points": key_points,

            "important_topics": important_topics,

            "created_at": datetime.utcnow()

        })
        track_library_file_event_by_filename(filename, "summary")

        return jsonify({

            "filename": filename,

            "subject": subject,

            "unit": unit,

            "summary": summary,

            "key_points": key_points,

            "important_topics": important_topics,

            "source": "gemini"

        })

    except Exception as e:
        return handle_exception(e, "Summary generation failed. Please try again.")


# ==========================================
# PART 2 STARTS HERE
# ==========================================
# ==========================================
# MCQ ROUTE
# ==========================================

@ai.route("/mcq/<path:filename>", methods=["GET"])
def generate_mcqs(filename):

    try:
        regenerate = request.args.get("regenerate") == "true"

        # -----------------------
        # Check MongoDB Cache
        # -----------------------

        existing = mcqs_collection.find_one({
            "filename": filename
        })

        if existing and not regenerate:

            return jsonify({

                "filename": filename,

                "mcqs": existing["mcqs"],

                "subject": existing.get("subject", ""),

                "unit": existing.get("unit", ""),

                "source": "cache"

            })

        if regenerate:
            mcqs_collection.delete_many({
                "filename": filename
            })

        study_text, summary_doc = study_text_for_filename(filename)

        if not study_text or not study_text.strip():
            return api_error(
                "Generate a summary first, or upload a readable PDF/PPTX/DOCX/TXT file.",
                400,
            )

        note = notes_collection.find_one({"filename": filename})
        subject = ""
        unit = ""
        if summary_doc:
            subject = summary_doc.get("subject", "")
            unit = summary_doc.get("unit", "")
        if note:
            subject = subject or note.get("subject", "")
            unit = unit or note.get("unit", "")

        prompt = f"""
You are an expert teacher.

Generate exactly 10 multiple choice questions.

Return ONLY valid JSON.

Format:

[
    {{
        "question":"Question",

        "options":[
            "Option A",
            "Option B",
            "Option C",
            "Option D"
        ],

        "answer":"Correct Option",
        "explanation":"One sentence explanation",
        "difficulty":"Easy | Medium | Hard",
        "topic":"Short topic name"
    }}
]

Rules:

1. JSON only.
2. No markdown.
3. Four options.
4. One correct answer.
5. Questions should help students prepare for exams.

Study Notes:

{study_text}
"""

        response = model.generate_content(prompt)

        mcqs = parse_mcqs(response.text)

        if not isinstance(mcqs, list) or len(mcqs) == 0:
            return api_error("AI did not return valid quiz questions. Please try again.", 502)

        mcqs_collection.insert_one({

            "filename": filename,

            "subject": subject,

            "unit": unit,

            "mcqs": mcqs,

            "created_at": datetime.utcnow()

        })
        track_library_file_event_by_filename(filename, "quiz")

        return jsonify({

            "filename": filename,

            "subject": subject,

            "unit": unit,

            "mcqs": mcqs,

            "source": "gemini"

        })

    except Exception as e:
        return handle_exception(e, "Quiz generation failed. Please try again.")


# ==========================================
# QUIZ SUBMIT
# ==========================================

@ai.route("/quiz/submit", methods=["POST"])
def submit_quiz():

    try:

        data = request.get_json(silent=True) or {}

        filename = data.get("filename")
        answers = data.get("answers", [])
        duration_seconds = data.get("duration_seconds")

        if not filename:
            return api_error("Filename is required", 400)

        if not isinstance(answers, list):
            return api_error("Answers must be a list", 400)

        mcq_doc = mcqs_collection.find_one({

            "filename": filename

        })

        if not mcq_doc:
            return api_error("MCQs not found. Generate a quiz first.", 404)

        mcqs = mcq_doc.get("mcqs", [])

        if not mcqs:
            return api_error("MCQs not found. Generate a quiz first.", 404)

        score = 0

        review = []

        for i, mcq in enumerate(mcqs):

            correct = mcq.get("answer") or mcq.get("correct_answer", "")

            user_answer = answers[i] if i < len(answers) else ""

            is_correct = (
                str(user_answer).strip().lower() == str(correct).strip().lower()
            )

            if is_correct:
                score += 1

            review.append({

                "question": mcq["question"],

                "selected": user_answer,

                "correct": correct,

                "is_correct": is_correct,

                "explanation": mcq.get("explanation", ""),

                "difficulty": mcq.get("difficulty", ""),

                "topic": mcq.get("topic", "")

            })

        total = len(mcqs)

        percentage = round(

            (score / total) * 100,

            2

        )

        if percentage >= 90:

            performance = "Excellent"

        elif percentage >= 75:

            performance = "Very Good"

        elif percentage >= 60:

            performance = "Good"

        elif percentage >= 40:

            performance = "Needs Improvement"

        else:

            performance = "Practice More"

        quiz_collection.insert_one({

            "filename": filename,

            "score": score,

            "total": total,

            "percentage": percentage,

            "performance": performance,

            "duration_seconds": duration_seconds,

            "submitted_at": datetime.utcnow()

        })

        return jsonify({

            "filename": filename,

            "score": score,

            "total": total,

            "percentage": percentage,

            "performance": performance,

            "review": review

        })

    except Exception as e:
        return handle_exception(e, "Unable to submit quiz results.")


# ==========================================
# QUIZ HISTORY
# ==========================================

@ai.route("/quiz/history", methods=["GET"])
def quiz_history():

    try:

        history = []

        quizzes = quiz_collection.find().sort([

            ("submitted_at", -1)

        ])

        for quiz in quizzes:

            history.append({

                "id": str(quiz["_id"]),

                "filename": quiz.get("filename"),

                "score": quiz.get("score"),

                "total": quiz.get("total"),

                "percentage": quiz.get("percentage"),

                "performance": quiz.get("performance"),

                "submitted_at": quiz.get("submitted_at").isoformat()

                if quiz.get("submitted_at") else ""

            })

        return jsonify(history)

    except Exception as e:
        return handle_exception(e, "Unable to load quiz history.")


# ==========================================
# CHAT WITH NOTES
# ==========================================

@ai.route("/chat/history/<path:filename>", methods=["GET"])
def chat_history(filename):

    try:

        messages = []

        for item in chat_collection.find({
            "filename": filename
        }).sort([
            ("created_at", 1)
        ]):

            messages.append({
                "question": item.get("question", ""),
                "answer": item.get("answer", ""),
                "created_at": item.get("created_at").isoformat()
                if item.get("created_at") else ""
            })

        return jsonify({
            "filename": filename,
            "messages": messages
        })

    except Exception as e:
        return handle_exception(e, "Unable to load chat history.")


@ai.route("/chat/<path:filename>", methods=["POST"])
def chat_with_note(filename):

    try:

        data = request.get_json() or {}
        question = data.get("question", "").strip()

        if not question:
            return jsonify({
                "message": "Question is required"
            }), 400

        summary_doc = summaries_collection.find_one({
            "filename": filename
        })

        note_text = ""

        if summary_doc:
            note_text = "\n".join(summary_doc.get("summary", []))
        else:
            note = notes_collection.find_one({"filename": filename})
            filepath, temp_path = material_filepath(note, filename)
            if not os.path.exists(filepath):
                return api_error("File not found", 404)
            try:
                note_text = extract_text(filepath)[:12000]
            finally:
                cleanup_temp_file(temp_path)

        prompt = f"""
You are EduVault AI, a helpful study tutor.

Answer the student's question using only the notes below.
If the answer is not present, say what is missing and suggest what to review.
Use concise markdown with bullets when helpful.

Student question:
{question}

Notes:
{note_text[:12000]}
"""

        response = model.generate_content(prompt)
        answer = response.text.strip()

        chat_collection.insert_one({
            "filename": filename,
            "question": question,
            "answer": answer,
            "created_at": datetime.utcnow()
        })
        track_library_file_event_by_filename(filename, "tutor")

        return jsonify({
            "filename": filename,
            "answer": answer
        })

    except Exception as e:
        return handle_exception(e, "Chat request failed. Please try again.")


# ==========================================
# DASHBOARD ANALYTICS
# ==========================================

@ai.route("/analytics", methods=["GET"])
def analytics():

    try:

        files_uploaded = notes_collection.count_documents({})
        summaries_generated = summaries_collection.count_documents({})
        mcqs_generated = mcqs_collection.count_documents({})
        quiz_attempts = quiz_collection.count_documents({})

        quiz_scores = list(quiz_collection.find({}, {"percentage": 1}))
        average_quiz_score = 0

        if quiz_scores:
            average_quiz_score = round(
                sum(float(q.get("percentage", 0)) for q in quiz_scores) / len(quiz_scores),
                1
            )

        recent_uploads = []

        for note in notes_collection.find().sort([
            ("created_at", -1)
        ]).limit(5):

            filename = note.get("filename", "")

            recent_uploads.append({
                "id": str(note["_id"]),
                "file_name": filename,
                "subject": note.get("subject", ""),
                "unit": note.get("unit", ""),
                "file_type": filename.split(".")[-1].upper()
                if "." in filename else "",
                "upload_date": note.get("created_at").isoformat()
                if note.get("created_at") else ""
            })

        activity = []
        today = datetime.utcnow().date()

        for offset in range(6, -1, -1):
            day = today - timedelta(days=offset)
            start = datetime.combine(day, datetime.min.time())
            end = start + timedelta(days=1)

            activity.append({
                "date": day.isoformat(),
                "uploads": notes_collection.count_documents({
                    "created_at": {
                        "$gte": start,
                        "$lt": end
                    }
                }),
                "quizzes": quiz_collection.count_documents({
                    "submitted_at": {
                        "$gte": start,
                        "$lt": end
                    }
                })
            })

        return jsonify({
            "files_uploaded": files_uploaded,
            "summaries_generated": summaries_generated,
            "mcqs_generated": mcqs_generated,
            "quiz_attempts": quiz_attempts,
            "average_quiz_score": average_quiz_score,
            "recent_uploads": recent_uploads,
            "activity": activity
        })

    except Exception as e:
        return handle_exception(e, "Unable to load analytics.")
